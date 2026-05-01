"""
從講義 ZIP 建 FAISS 向量庫並打包成 RAG ZIP；含 process_zip_to_docs（與建庫相同副檔名規則）。
Embeddings 由呼叫端傳入 API key，不從環境變數讀取。

unit_type（與 Rag_Unit.unit_type／POST /rag/tab/build-rag-zip 之 unit_types 對齊）：
  0／1：Office／PDF＋Markdown
  2／4：僅 Markdown（.md）
  3：僅 MP3（逐字稿依 Deepgram，需設定 Deepgram API Key）
"""

import io
import os
import shutil
import tempfile
import zipfile
from pathlib import Path

from langchain_community.document_loaders import Docx2txtLoader, PyPDFLoader
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document
from langchain_openai import OpenAIEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter

from utils.zip_utils import fix_encoding

# unit_type 常數（與 routers/zip.py 及 Rag_Unit.unit_type 一致）
UNIT_TYPE_DEFAULT = 0
UNIT_TYPE_RAG = 1
UNIT_TYPE_TEXT = 2
UNIT_TYPE_MP3 = 3
UNIT_TYPE_YOUTUBE = 4

# 支援載入的文件副檔名（小寫）；.doc/.ppt 依賴 unstructured
_SUPPORTED_DOC_EXTS: frozenset[str] = frozenset(
    {".pdf", ".doc", ".docx", ".ppt", ".pptx", ".md"}
)


def _documents_from_pptx(path: Path) -> list[Document]:
    """將 .pptx 各投影片文字合併為單一 Document（slide N\n<text>）。"""
    from pptx import Presentation

    prs = Presentation(str(path))
    chunks: list[str] = []
    for idx, slide in enumerate(prs.slides):
        lines: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            t = (shape.text or "").strip()
            if t:
                lines.append(t)
        if lines:
            chunks.append(f"Slide {idx + 1}\n" + "\n".join(lines))
    text = "\n\n".join(chunks).strip()
    if not text:
        return []
    return [Document(page_content=text, metadata={"source": str(path)})]


def _documents_from_md(path: Path) -> list[Document]:
    """讀取 .md 檔案原文（保留 Markdown 語法）為 Document。"""
    try:
        raw = path.read_bytes()
        try:
            text = raw.decode("utf-8")
        except UnicodeDecodeError:
            text = raw.decode("utf-8", errors="replace")
    except OSError:
        return []
    text = text.strip()
    if not text:
        return []
    return [Document(page_content=text, metadata={"source": str(path)})]


def _load_docs_from_file(path: Path) -> list[Document]:
    """依副檔名選用對應 loader 載入單一檔案；不支援的副檔名回傳空 list。"""
    ext = path.suffix.lower()
    if ext == ".md":
        return _documents_from_md(path)
    if ext not in _SUPPORTED_DOC_EXTS:
        return []
    try:
        if ext == ".pdf":
            return PyPDFLoader(str(path)).load()
        if ext == ".docx":
            return Docx2txtLoader(str(path)).load()
        if ext == ".doc":
            from langchain_community.document_loaders import UnstructuredWordDocumentLoader
            return UnstructuredWordDocumentLoader(str(path)).load()
        if ext == ".pptx":
            return _documents_from_pptx(path)
        if ext == ".ppt":
            from langchain_community.document_loaders import UnstructuredPowerPointLoader
            return UnstructuredPowerPointLoader(str(path)).load()
    except Exception:
        return []
    return []


def _is_skipped_extract_path(path: Path) -> bool:
    """過濾 macOS 產生的隱藏項目（__MACOSX、.DS_Store）。"""
    s = str(path)
    return "__MACOSX" in s or path.name == ".DS_Store"


def _collect_documents_after_extract(extract_dir: Path, unit_type: int) -> list[Document]:
    """依 unit_type 從已解壓目錄收集 Document（須先完成 ZIP 解壓）。"""
    ut = unit_type
    if ut < 0 or ut > 4:
        ut = UNIT_TYPE_DEFAULT

    # unit_type 2（文字）/ 4（YouTube）：僅讀 .md 檔
    if ut in (UNIT_TYPE_TEXT, UNIT_TYPE_YOUTUBE):
        md_paths: list[Path] = []
        for path in extract_dir.rglob("*"):
            if not path.is_file() or _is_skipped_extract_path(path):
                continue
            if path.suffix.lower() == ".md":
                md_paths.append(path)
        md_paths.sort(key=lambda p: str(p))
        out: list[Document] = []
        for p in md_paths:
            out.extend(_documents_from_md(p))
        return out

    # unit_type 3（MP3）：取第一個音訊檔，以 Deepgram 轉逐字稿
    if ut == UNIT_TYPE_MP3:
        from utils.media_transcript import transcribe_audio_bytes_deepgram

        mp3_paths: list[Path] = []
        for path in extract_dir.rglob("*"):
            if not path.is_file() or _is_skipped_extract_path(path):
                continue
            if path.suffix.lower() == ".mp3":
                mp3_paths.append(path)
        mp3_paths.sort(key=lambda p: str(p))
        if not mp3_paths:
            return []
        p0 = mp3_paths[0]
        data = p0.read_bytes()
        text, _elapsed = transcribe_audio_bytes_deepgram(data, suffix=".mp3")
        text = (text or "").strip()
        if not text:
            return []
        return [Document(page_content=text, metadata={"source": str(p0)})]

    # unit_type 0／1：Office＋Markdown，全部支援格式皆嘗試載入
    all_docs: list[Document] = []
    for path in extract_dir.rglob("*"):
        if not path.is_file() or _is_skipped_extract_path(path):
            continue
        all_docs.extend(_load_docs_from_file(path))
    return all_docs


def _empty_docs_user_message(unit_type: int) -> str:
    """依 unit_type 回傳無文件時的友善錯誤訊息。"""
    ut = unit_type if 0 <= unit_type <= 4 else UNIT_TYPE_DEFAULT
    if ut in (UNIT_TYPE_TEXT, UNIT_TYPE_YOUTUBE):
        return "ZIP 內無可讀 .md（文字／YouTube 單元僅使用 Markdown）"
    if ut == UNIT_TYPE_MP3:
        return (
            "ZIP 內無 .mp3，或 Deepgram 逐字稿為空。聽力單元須含音檔；並請設定 Deepgram API Key "
            "（環境變數 DEEPGRAM_API_KEY）"
        )
    return "ZIP 內無可讀文件（支援：.pdf .doc .docx .ppt .pptx .md）"


def process_zip_to_docs(
    zip_path: Path, extract_dir: Path, unit_type: int = UNIT_TYPE_DEFAULT
) -> list[Document]:
    """
    解壓 ZIP 到 extract_dir，依 unit_type 載入檔案為 Document 列表。

    - 0／1：.pdf, .doc, .docx, .ppt, .pptx, .md
    - 2／4：僅 .md
    - 3：僅 .mp3（Deepgram 轉逐字稿後嵌入）

    自動過濾 __MACOSX、.DS_Store，並修正 cp437/big5 編碼。
    """
    extract_dir = Path(extract_dir)
    extract_dir.mkdir(parents=True, exist_ok=True)

    with zipfile.ZipFile(zip_path, "r") as z:
        for raw_name in z.namelist():
            if raw_name.endswith("/"):
                continue
            decoded = fix_encoding(raw_name)
            if "__MACOSX" in decoded or ".DS_Store" in decoded:
                continue
            # 防止路徑穿越：將 .. 替換為 _
            safe_path = extract_dir / decoded.replace("..", "_")
            safe_path.parent.mkdir(parents=True, exist_ok=True)
            safe_path.write_bytes(z.read(raw_name))

    return _collect_documents_after_extract(extract_dir, unit_type)


def build_faiss_zip_from_docs(
    documents: list[Document],
    api_key: str,
    chunk_size: int,
    chunk_overlap: int,
) -> bytes:
    """
    將 Document 列表做切分、Embedding、FAISS 建索引，打包成 ZIP bytes 回傳。

    Embedding 模型固定為 text-embedding-3-small（須與出題／批改路徑一致，維持索引維度相容）。
    """
    if not documents:
        raise ValueError("無文件可處理")

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
    )
    split_docs = text_splitter.split_documents(documents)

    embeddings = OpenAIEmbeddings(model="text-embedding-3-small", api_key=api_key)
    vectorstore = FAISS.from_documents(split_docs, embeddings)

    tmpdir = Path(tempfile.mkdtemp())
    try:
        vectorstore.save_local(str(tmpdir))
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
            for root, _dirs, files in os.walk(tmpdir):
                for f in files:
                    fp = Path(root) / f
                    arcname = os.path.relpath(fp, tmpdir)
                    zf.write(fp, arcname)
        return buf.getvalue()
    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)


def make_rag_zip_from_zip_path(
    zip_path: Path,
    api_key: str,
    chunk_size: int,
    chunk_overlap: int,
    unit_type: int = UNIT_TYPE_DEFAULT,
) -> bytes:
    """
    從一個 ZIP 路徑：解壓 → 載入文件 → 切分 → Embedding → FAISS → 打包成 ZIP bytes。

    unit_type 控制載入邏輯（見模組說明）；與 build-rag-zip 各任務之 unit_types 分段對齊。
    """
    tmp_extract = Path(tempfile.mkdtemp())
    try:
        all_docs = process_zip_to_docs(zip_path, tmp_extract, unit_type=unit_type)
        if not all_docs:
            raise ValueError(_empty_docs_user_message(unit_type))
        return build_faiss_zip_from_docs(
            all_docs, api_key, chunk_size=chunk_size, chunk_overlap=chunk_overlap
        )
    finally:
        shutil.rmtree(tmp_extract, ignore_errors=True)
